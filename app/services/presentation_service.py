"""
presentation_service.py

Generates a dynamically-themed 12-slide PPTX marketing strategy deck.

Visual language (inspired by reference consulting decks):
  - Card-based layouts with colored left-border or top-border accents
  - Consistent slide header: title (left) + brand name (right) + full-width rule
  - Multi-column grids: 2×2 for market/objectives, 3-col for audience/KPIs
  - Typography hierarchy: LABEL → Bold Header → Regular Body within each card
  - Campaign examples slide referencing real analogous companies

python-pptx 1.0.2: solid fills only (no native gradients).
Canvas: 10" × 5.625" widescreen (9,144,000 × 5,143,500 EMU).
"""

import json
import logging
import os
import random
import traceback

from openai import OpenAI

from app.core.config import settings

logger = logging.getLogger(__name__)

SLIDE_W = 9_144_000
SLIDE_H = 5_143_500

# ---------------------------------------------------------------------------
# Template definitions
# Each template file lives in "example pptx/" at the project root.
# TEMPLATE_MAP: key → (filename, is_dark_background)
# ---------------------------------------------------------------------------
_EXAMPLES_DIR_RELATIVE = os.path.join(os.path.dirname(__file__), "..", "..", "example pptx")
_EXAMPLES_DIR_CWD = os.path.join(os.getcwd(), "example pptx")
_EXAMPLES_DIR = (
    _EXAMPLES_DIR_RELATIVE
    if os.path.isdir(_EXAMPLES_DIR_RELATIVE)
    else _EXAMPLES_DIR_CWD
)

TEMPLATE_MAP = {
    "corporate": (
        "Black and Red Elegant Corporate Marketing Plan Presentation.pptx",
        True,        # dark background
        "diagonal",  # matching design_variant
    ),
    "digital": (
        "Cream Blue Creative Modern Digital Marketing Presentation (1).pptx",
        False,
        "bold_contrast",
    ),
    "organic": (
        "Green And White Illustrative Marketing Plan Presentation.pptx",
        False,
        "top_band",
    ),
    "creative": (
        "Orange and Cream Illustration Marketing Plan Presentation.pptx",
        False,
        "minimal",
    ),
}

# Industry keyword → template key
_TEMPLATE_INDUSTRY: dict[str, str] = {
    "finance": "corporate", "fintech": "corporate", "banking": "corporate",
    "b2b": "corporate", "consulting": "corporate", "insurance": "corporate",
    "legal": "corporate",
    "tech": "digital", "saas": "digital", "software": "digital",
    "startup": "digital", "ai": "digital", "media": "digital", "gaming": "digital",
    "wellness": "organic", "health": "organic", "food": "organic",
    "sustainability": "organic", "ecommerce": "organic", "retail": "organic",
    "fitness": "organic", "restaurant": "organic",
    "fashion": "creative", "beauty": "creative", "luxury": "creative",
    "travel": "creative", "education": "creative", "real estate": "creative",
    "lifestyle": "creative",
}

# Tone keyword → template key (applied after industry, so tone can override)
_TEMPLATE_TONE: dict[str, str] = {
    "corporate": "corporate", "professional": "corporate", "trustworthy": "corporate",
    "elegant": "corporate", "sophisticated": "corporate",
    "innovative": "digital", "modern": "digital", "bold": "digital",
    "edgy": "digital", "dynamic": "digital",
    "clean": "organic", "minimal": "organic", "organic": "organic",
    "fresh": "organic", "sustainable": "organic",
    "playful": "creative", "fun": "creative", "energetic": "creative",
    "youthful": "creative", "vibrant": "creative", "creative": "creative",
}

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

RECT       = MSO_AUTO_SHAPE_TYPE.RECTANGLE
ROUND_RECT = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE


class PresentationService:
    def __init__(self):
        self.client = OpenAI(api_key=settings.OPENAI_API_KEY)
        self.model = settings.GPT_MODEL

    # =========================================================================
    # PUBLIC: structure_content
    # =========================================================================
    def structure_content(self, questionnaire: dict, analysis: dict) -> dict:
        """
        Generates text content for 12 slides. Returns {"slides": [...]} JSON.
        """
        meta      = questionnaire.get("project_metadata", {}) or {}
        creative  = questionnaire.get("the_creative_goal", {}) or {}
        market    = questionnaire.get("market_context", {}) or {}
        product   = questionnaire.get("product_definition", {}) or {}
        audience  = questionnaire.get("target_audience", {}) or {}

        brand_name    = meta.get("brand_name", "Brand")
        industry      = meta.get("industry", "")
        country       = meta.get("target_country", "")
        website       = meta.get("website_url", "")
        channels      = ", ".join(creative.get("specific_channels", []) or [])
        usp           = product.get("unique_selling_proposition", "")
        competitors   = ", ".join(market.get("main_competitors", []) or [])
        tone          = creative.get("desired_tone_of_voice", "")
        objectives    = ", ".join(creative.get("marketing_objectives", []) or [])
        audience_desc = audience.get("description", "") or audience.get("primary_segment", "")
        objections    = ", ".join(market.get("main_objections", []) or [])

        hooks       = analysis.get("hooks", [])
        angles      = analysis.get("angles", [])
        pivot       = analysis.get("creative_pivot", "")
        awareness   = analysis.get("brand_awareness_strategy", {}) or {}
        channel_tac = awareness.get("channel_tactics", [])
        quick_wins  = awareness.get("quick_wins", [])
        positioning = awareness.get("positioning_recommendation", "")

        perplexity_snapshot  = json.dumps(analysis.get("perplexity_research_snapshot", {}), indent=2)[:3000]
        brand_audit_snapshot = json.dumps(analysis.get("brand_audit_snapshot", {}), indent=2)[:1500]
        news_snapshot        = json.dumps(analysis.get("news_snapshot", {}), indent=2)[:2000]

        system_prompt = (
            "You are a Senior Marketing Strategist at a top-tier creative agency. "
            "You build board-ready pitch decks with deep strategic insight — "
            "specific competitor names, real benchmarks, named techniques. "
            "Output must be valid JSON only. No markdown, no explanations."
        )

        user_content = f"""
# Brand Briefing
Brand: {brand_name} | Industry: {industry} | Country: {country}
Website: {website}
USP: {usp}
Target Audience: {audience_desc}
Channels: {channels} | Tone: {tone}
Objectives: {objectives}
Competitors: {competitors}
Objections: {objections}

# Strategy Consensus
Creative Pivot: {pivot}
Top Angles: {json.dumps(angles)}
Top Hooks: {json.dumps(hooks[:7])}
Channel Tactics: {json.dumps(channel_tac)}
Quick Wins: {json.dumps(quick_wins)}
Positioning: {positioning}

# Research Data
## Competitor & Brand Awareness (Perplexity):
{perplexity_snapshot}

## Brand Website Audit:
{brand_audit_snapshot}

## Press & News Coverage (NewsAPI):
{news_snapshot}

# TASK
Create exactly 16 slides. Return a JSON object with a "slides" array.
Be SPECIFIC — cite real competitor names, research findings, real benchmarks.
Use news coverage and competitor press activity to ground claims in real evidence.

## SLIDE 1 — type: "title"
- title: Punchy campaign tagline ≤8 words capturing the strategic angle
- subtitle: One crisp brand promise ≤15 words

## SLIDE 2 — type: "company_intro"  (CONCISE company overview)
- title: "About {brand_name}"
- headline: Bold strategic statement ≤10 words (the core value in one line)
- description: Exactly 2 sentences describing what {brand_name} does and for whom
- kvp: Array of 3 objects: {{"label": "Short Label ≤3 words", "description": "≤8 words"}}
  Cover the 3 core value pillars (product strength, user benefit, market edge)

## SLIDE 3 — type: "two_by_two"  (Market overview)
- title: "The Market"
- cards: Array of exactly 4 objects: {{"label": "LABEL", "header": "Short bold title", "body": "1-2 sentences"}}
  Cards: VALIDATION (funding/backing/traction), TRACTION (current user/revenue state),
         CHALLENGE (the main market problem with stat from research), OPPORTUNITY (the gap to fill)

## SLIDE 4 — type: "single_card"  (The one core challenge)
- title: "The Main Challenge"
- label: Context label (e.g. "Expansion Phase", "Growth Barrier")
- headline: Bold challenge name ≤6 words
- body: 2 sentences MAX — specific to {brand_name}'s situation. Reference real data from research.

## SLIDE 5 — type: "persona_detail"  (Primary Persona — named profile)
- title: "Target Audience"
- subtitle: Persona category ≤5 words (e.g. "Growth-Stage Marketing Leaders")
- name: Real first name + age (e.g. "Maya, 32")
- role: Job title (e.g. "VP of Marketing")
- company: Company type/sector (e.g. "Series B SaaS Startup")
- tags: Array of 2-3 short role-descriptor tags (e.g. ["Strategic", "B2B"])
- quote: First-person quote ≤25 words capturing their core frustration, in their voice
- cards: Array of exactly 4 objects:
  [{{"label": "Scope of Responsibility", "body": "1-2 sentences specific to this persona"}},
   {{"label": "Primary Concerns", "body": "1-2 sentences specific to this persona"}},
   {{"label": "Core Challenge", "body": "1-2 sentences specific to this persona"}},
   {{"label": "Critical Need", "body": "1-2 sentences specific to this persona"}}]

## SLIDE 6 — type: "persona_detail"  (Secondary Persona — named profile)
- Same structure as Slide 5 but DIFFERENT persona (different role/seniority/demographic)

## SLIDE 7 — type: "three_col"  (Audience deep-dive)
- title: "Audience Deep-Dive"
- columns: Array of exactly 3 objects:
  [{{"label": "DEMOGRAPHICS", "header": "Who They Are", "items": ["age range", "location", "income/role", "education"]}},
   {{"label": "PSYCHOGRAPHICS", "header": "How They Think", "items": ["core value 1", "core value 2", "media habits", "buying trigger"]}},
   {{"label": "KEY BUYERS", "header": "Decision Makers", "items": ["Title 1", "Title 2", "Title 3", "Title 4"]}}]

## SLIDE 8 — type: "two_col"  (Competitive landscape)
- title: "Competitive Landscape"
- left: {{"label": "COMPETITOR GAPS", "header": "What They Miss", "items": ["gap 1 (name the competitor)", "gap 2 (name the competitor)", "gap 3"]}}
- right: {{"label": "OUR EDGE", "header": "{brand_name}'s Advantage", "items": ["advantage 1", "advantage 2", "advantage 3"]}}

## SLIDE 9 — type: "three_col"  (3V's of Brand Identity)
- title: "Brand Identity — The 3 V's"
- columns: exactly 3 objects:
  [{{"label": "VISION", "header": "Where we're going", "items": ["The world we're building toward", "The change we want to create", "Our north star metric"]}},
   {{"label": "VALUES", "header": "What we stand for", "items": ["Core value 1", "Core value 2", "Core value 3"]}},
   {{"label": "VOICE", "header": "How we sound", "items": ["Tone descriptor 1", "Tone descriptor 2", "What we never say", "Platform where voice shines"]}}]

## SLIDE 10 — type: "three_col"  (Campaign barriers)
- title: "Barriers to Success"
- columns: exactly 3 objects, one per barrier type:
  [{{"label": "MARKET BARRIER", "header": "Adoption & Trust Gap",
    "items": [
      "Specific adoption or trust challenge drawn from news/community research — name the root cause",
      "A real data point or user sentiment that proves this barrier exists",
      "Why standard marketing approaches fail to overcome it for {brand_name}"
    ]}},
   {{"label": "COMPETITIVE BARRIER", "header": "Competitor Inertia",
    "items": [
      "Name the specific competitor(s) creating this inertia and what advantage they hold",
      "The switching cost or lock-in mechanism that keeps users from moving to {brand_name}",
      "What {brand_name} must do differently to overcome this competitor advantage"
    ]}},
   {{"label": "EXECUTION BARRIER", "header": "Internal Challenge",
    "items": [
      "The internal capability or perception gap {brand_name} must close",
      "Why this is specifically hard given {brand_name}'s current market position or brand history",
      "The one lever {brand_name} can pull to address this from the inside out"
    ]}}]

## SLIDE 11 — type: "campaign"
- title: "Campaign: [TECHNIQUE]" — pick from: Viral Referral Loop / UGC Flywheel / Micro-Influencer Tier / Community-Led Growth / Social Proof Cascade / Challenge Campaign / Waitlist Launch / Paid-Organic Flywheel / Content-Led SEO / Partnership Co-Marketing
- subtitle: One-line campaign concept name (e.g. "The Trust Engine")
- content: 3 bullets:
  (a) How this technique works specifically for {brand_name}
  (b) The trigger or hook that activates it — reference a real pain point from competitor research or news
  (c) Specific success metric with a number (e.g. "Target 25% referral-driven signups in 60 days")

## SLIDE 12 — type: "campaign"
- Same structure as Slide 11 but DIFFERENT technique and DIFFERENT channel focus

## SLIDE 13 — type: "campaign_examples"
- title: "Campaign Inspiration"
- examples: Array of exactly 2 objects — REAL companies with analogous campaigns:
  {{"company": "Real Company Name", "technique": "Short Technique Label", "strategy": "Strategy Theme Name", "items": ["point 1", "point 2", "point 3"]}}
  Pick companies like Canva, Gong, Notion, Figma, HubSpot, Slack, Drift, Airbnb, Duolingo, Dropbox
  whose campaigns are directly analogous to what we're recommending for {brand_name}.

## SLIDE 14 — type: "hooks"
- title: "Marketing Hooks"
- content: Exactly 5 hooks. Each ≤20 words. Ground at least 2 in findings from news or competitor research.

## SLIDE 15 — type: "content"  (Cross-platform plan)
- title: "Cross-Platform Plan"
- content: 3 bullets, each covering a different platform/channel:
  (a) Primary channel: why chosen, what content format, KPI
  (b) Secondary channel: why chosen, what content format, KPI
  (c) Supporting channel: why chosen, what content format, KPI

## SLIDE 16 — type: "kpis"
- title: "KPI's"
- columns: Array of exactly 3 objects:
  [{{"label": "AWARENESS", "subtitle": "Reach & Visibility", "metrics": [{{"number": "X", "desc": "metric name (source)"}}]}},
   {{"label": "ENGAGEMENT", "subtitle": "Interaction & Interest", "metrics": [...]}},
   {{"label": "CONVERSION", "subtitle": "Leads & Action", "metrics": [...]}}]
  Each column has 2-3 metrics. Numbers must be SPECIFIC with industry benchmark sources.
  Example: "8.5%" with desc "Avg engagement rate (HubSpot 2024 benchmark: 3.2%)"

RULES:
- Slide 2: description MUST be exactly 2 sentences. No more.
- Slide 4: body MUST be ≤2 sentences, specific to {brand_name}.
- Slides 5 & 6: names MUST be real first names; quote MUST be in their voice; cards must be role-specific, not generic.
- Slide 9: 3V's content must be specific to {brand_name}, not generic platitudes.
- Slide 13: companies must be REAL. Technique must be analogous to slides 11-12.
- All text: concise, specific. No filler phrases like "leverage synergies".
- Return ONLY the JSON object. No markdown wrapping.
"""

        try:
            response = self.client.chat.completions.create(
                model=self.model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_content},
                ],
                response_format={"type": "json_object"},
                temperature=0.7,
            )
            return json.loads(response.choices[0].message.content)
        except Exception as e:
            logger.error(f"Presentation structuring error: {e}")
            return {"error": str(e), "slides": []}

    # =========================================================================
    # PUBLIC: generate_pptx
    # =========================================================================
    def generate_pptx(
        self,
        slides_data: dict,
        output_path: str,
        questionnaire: dict = None,
    ) -> str:
        from pptx import Presentation
        from pptx.util import Emu

        try:
            theme = self._derive_theme(questionnaire or {})

            template_path = theme.get("template_path")
            if template_path and os.path.isfile(template_path):
                prs = Presentation(template_path)
                logger.info(f"Using template: {os.path.basename(template_path)}")
                # Remove the template's existing content slides (keep slide master/layouts)
                _NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
                xml_slide_list = prs.slides._sldIdLst
                for sld_id_el in list(xml_slide_list):
                    rId = sld_id_el.get(_NS_R + "id")
                    if rId:
                        prs.part.drop_rel(rId)
                    xml_slide_list.remove(sld_id_el)
            else:
                prs = Presentation()
            prs.slide_width = Emu(SLIDE_W)
            prs.slide_height = Emu(SLIDE_H)

            blank_layout = (
                prs.slide_layouts[6]
                if len(prs.slide_layouts) > 6
                else prs.slide_layouts[0]
            )

            slides = slides_data.get("slides", [])
            total = len(slides)
            for idx, slide_info in enumerate(slides):
                slide = prs.slides.add_slide(blank_layout)
                stype = slide_info.get("type", "content")
                num = idx + 1

                if stype == "title":
                    self._build_slide_title(slide, slide_info, theme, num, total)
                elif stype == "company_intro":
                    self._build_slide_company_intro(slide, slide_info, theme, num, total)
                elif stype == "two_by_two":
                    self._build_slide_two_by_two(slide, slide_info, theme, num, total)
                elif stype == "single_card":
                    self._build_slide_single_card(slide, slide_info, theme, num, total)
                elif stype == "three_col":
                    self._build_slide_three_col(slide, slide_info, theme, num, total)
                elif stype == "two_col":
                    self._build_slide_two_col(slide, slide_info, theme, num, total)
                elif stype == "persona_detail":
                    self._build_slide_persona_detail(slide, slide_info, theme, num, total)
                elif stype == "campaign":
                    self._build_slide_campaign(slide, slide_info, theme, num, total)
                elif stype == "campaign_examples":
                    self._build_slide_campaign_examples(slide, slide_info, theme, num, total)
                elif stype == "hooks":
                    self._build_slide_hooks(slide, slide_info, theme, num, total)
                elif stype == "kpis":
                    self._build_slide_kpis(slide, slide_info, theme, num, total)
                elif stype in ("roadmap", "next_steps"):
                    self._build_slide_roadmap(slide, slide_info, theme, num, total)
                else:
                    # "content", "social", and fallback
                    self._build_slide_content(slide, slide_info, theme, num, total)

            prs.save(output_path)
            logger.info(f"PPTX saved to {output_path}")
            return output_path

        except Exception as e:
            logger.exception("PPTX generation error")
            raise RuntimeError(f"PPTX generation failed: {e}") from e

    # =========================================================================
    # PRIVATE: _derive_theme
    # =========================================================================
    def _derive_theme(self, questionnaire: dict) -> dict:
        from pptx.dml.color import RGBColor

        meta     = questionnaire.get("project_metadata", {}) or {}
        creative = questionnaire.get("the_creative_goal", {}) or {}

        industry_raw = (meta.get("industry", "") or "").lower()
        tone_raw     = (creative.get("desired_tone_of_voice", "") or "").lower()
        brand_name   = (meta.get("brand_name", "Brand") or "Brand").strip()
        website_url  = str(meta.get("website_url", "") or "")
        country      = (meta.get("target_country", "") or "").strip()
        channels     = creative.get("specific_channels", []) or []

        PALETTES = {
            "tech":           ("3b82f6", "6366f1", "06b6d4", "0f172a"),
            "saas":           ("3b82f6", "6366f1", "06b6d4", "0f172a"),
            "software":       ("3b82f6", "6366f1", "06b6d4", "0f172a"),
            "startup":        ("6366f1", "8b5cf6", "22d3ee", "0f172a"),
            "ai":             ("6366f1", "8b5cf6", "22d3ee", "0f172a"),
            "wellness":       ("10b981", "14b8a6", "84cc16", "f0fdf4"),
            "health":         ("10b981", "14b8a6", "84cc16", "f0fdf4"),
            "fitness":        ("10b981", "14b8a6", "a3e635", "f0fdf4"),
            "beauty":         ("a855f7", "ec4899", "f59e0b", "1a0a2e"),
            "fashion":        ("a855f7", "ec4899", "f59e0b", "1a0a2e"),
            "luxury":         ("a855f7", "c084fc", "f59e0b", "0d0014"),
            "finance":        ("1e40af", "0ea5e9", "f59e0b", "0c1a3a"),
            "fintech":        ("1e40af", "0ea5e9", "f59e0b", "0c1a3a"),
            "banking":        ("1e40af", "0369a1", "f59e0b", "0c1a3a"),
            "b2b":            ("1e40af", "0ea5e9", "d97706", "0c1a3a"),
            "consulting":     ("1e40af", "0ea5e9", "d97706", "0c1a3a"),
            "food":           ("f97316", "ef4444", "fbbf24", "fff7ed"),
            "restaurant":     ("f97316", "ef4444", "fbbf24", "fff7ed"),
            "sustainability": ("10b981", "059669", "84cc16", "f0fdf4"),
            "education":      ("0284c7", "0ea5e9", "f59e0b", "eff6ff"),
            "ecommerce":      ("e11d48", "f43f5e", "f59e0b", "fff1f2"),
            "retail":         ("e11d48", "f43f5e", "f59e0b", "fff1f2"),
            "media":          ("7c3aed", "8b5cf6", "f43f5e", "0f0720"),
            "gaming":         ("7c3aed", "8b5cf6", "22d3ee", "0f0720"),
            "travel":         ("0891b2", "06b6d4", "f59e0b", "ecfeff"),
            "real estate":    ("059669", "10b981", "f59e0b", "f0fdf4"),
        }
        DEFAULT = ("3b82f6", "8b5cf6", "06b6d4", "0f172a")

        palette = DEFAULT
        for keyword, pal in PALETTES.items():
            if keyword in industry_raw:
                palette = pal
                break

        primary_hex, secondary_hex, accent_hex, bg_hex = palette

        tone_tokens = set(t.strip() for t in tone_raw.replace(",", " ").split())
        force_dark    = bool(tone_tokens & {"bold", "innovative", "edgy", "premium", "luxury", "dramatic"})
        force_light   = bool(tone_tokens & {"professional", "trustworthy", "clean", "minimal", "corporate"})
        force_gold    = bool(tone_tokens & {"premium", "luxury", "elegant", "exclusive"})
        force_vibrant = bool(tone_tokens & {"playful", "fun", "energetic", "youthful", "vibrant"})

        is_light_bg = bg_hex[0] in ("f", "e")
        if force_dark and is_light_bg:
            bg_hex = "0f172a"
            is_light_bg = False
        if force_light and not is_light_bg:
            bg_hex = "eee9e2"
            is_light_bg = True
        if force_gold:
            accent_hex = "f59e0b"
        if force_vibrant:
            accent_hex = secondary_hex

        text_light = "f1f5f9"
        text_dark  = "0f172a"
        muted      = "94a3b8" if not is_light_bg else "64748b"
        text_main  = text_dark if is_light_bg else text_light

        # Card background: slightly lighter/darker than slide bg
        if not is_light_bg:
            # Slightly lighter dark
            bg_int = int(bg_hex, 16)
            card_r = min(255, ((bg_int >> 16) & 0xFF) + 20)
            card_g = min(255, ((bg_int >> 8) & 0xFF) + 20)
            card_b = min(255, (bg_int & 0xFF) + 30)
            card_bg_hex = f"{card_r:02x}{card_g:02x}{card_b:02x}"
        else:
            card_bg_hex = "dbd5cd"

        def h(hex_str: str):
            s = hex_str.lstrip("#")
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

        # ── Template + design_variant: selected deterministically from input ──
        # Industry keyword takes priority; tone can override.
        template_key = "digital"  # default
        for kw, tkey in _TEMPLATE_INDUSTRY.items():
            if kw in industry_raw:
                template_key = tkey
                break
        for kw in tone_tokens:
            if kw in _TEMPLATE_TONE:
                template_key = _TEMPLATE_TONE[kw]
                break

        template_fname, template_is_dark, design_variant = TEMPLATE_MAP[template_key]
        template_path = os.path.join(_EXAMPLES_DIR, template_fname)
        use_template = os.path.isfile(template_path)

        # If using a template, override dark_theme to match the template's bg
        if use_template:
            is_light_bg = not template_is_dark
            text_main = text_dark if is_light_bg else text_light
            muted = "94a3b8" if not is_light_bg else "64748b"

        return {
            "primary":          h(primary_hex),
            "secondary":        h(secondary_hex),
            "accent":           h(accent_hex),
            "bg":               h(bg_hex),
            "card_bg":          h(card_bg_hex),
            "text_main":        h(text_main),
            "text_light":       h(text_light),
            "text_dark":        h(text_dark),
            "muted":            h(muted),
            "dark_theme":       not is_light_bg,
            "design_variant":   design_variant,
            "brand_name":       brand_name,
            "website_url":      website_url,
            "target_country":   country,
            "channels":         channels,
            "template_path":    template_path if use_template else None,
            "use_template_bg":  use_template,
        }

    # =========================================================================
    # PRIVATE: low-level drawing utilities
    # =========================================================================

    def _set_background(self, slide, color, theme=None) -> None:
        if theme and theme.get("use_template_bg"):
            return  # let template master background show through
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_bg_accent(self, slide, theme) -> None:
        """
        Geometric depth element — shape varies by design_variant:
          diagonal / bold_contrast → diagonal freeform triangle (bottom-right)
          top_band                 → bottom-left triangle
          minimal                  → very subtle rect strip
        """
        if theme.get("use_template_bg"):
            return  # template has its own decorative background elements
        from pptx.dml.color import RGBColor
        bg = theme["bg"]
        r, g, b = bg[0], bg[1], bg[2]
        if theme["dark_theme"]:
            nr = min(255, r + 20); ng = min(255, g + 20); nb = min(255, b + 28)
        else:
            nr = max(0, r - 10);   ng = max(0, g - 10);   nb = max(0, b - 10)
        depth_color = RGBColor(nr, ng, nb)
        variant = theme.get("design_variant", "diagonal")

        try:
            if variant == "top_band":
                # Bottom-left triangle
                builder = slide.shapes.build_freeform(0, SLIDE_H)
                builder.add_line_to(int(SLIDE_W * 0.48), SLIDE_H)
                builder.add_line_to(0, int(SLIDE_H * 0.62))
                tri = builder.convert_to_shape()
            elif variant == "minimal":
                # Thin right-edge strip instead
                self._add_rect(slide, SLIDE_W - int(SLIDE_W * 0.018), 0,
                                int(SLIDE_W * 0.018), SLIDE_H, depth_color)
                return
            else:
                # diagonal / bold_contrast: bottom-right triangle
                builder = slide.shapes.build_freeform(SLIDE_W, SLIDE_H)
                builder.add_line_to(SLIDE_W, int(SLIDE_H * 0.28))
                builder.add_line_to(int(SLIDE_W * 0.50), SLIDE_H)
                tri = builder.convert_to_shape()
            tri.fill.solid()
            tri.fill.fore_color.rgb = depth_color
            tri.line.fill.background()
        except Exception:
            w = int(SLIDE_W * 0.42)
            h = int(SLIDE_H * 0.55)
            self._add_rect(slide, SLIDE_W - w, SLIDE_H - h, w, h, depth_color)

    def _add_slide_counter(self, slide, num: int, total: int, theme: dict) -> None:
        """Full-width footer bar: brand name (left), website (center), slide counter (right)."""
        from pptx.dml.color import RGBColor
        footer_h = int(SLIDE_H * 0.060)
        footer_t = SLIDE_H - footer_h

        bg = theme["bg"]
        if theme["dark_theme"]:
            nr = min(255, bg[0] + 16)
            ng = min(255, bg[1] + 16)
            nb = min(255, bg[2] + 22)
        else:
            nr = max(0, bg[0] - 13)
            ng = max(0, bg[1] - 13)
            nb = max(0, bg[2] - 13)
        footer_color = RGBColor(nr, ng, nb)

        self._add_rect(slide, 0, footer_t, SLIDE_W, footer_h, footer_color)
        # Thin accent rule above footer
        self._add_rect(slide, 0, footer_t, SLIDE_W, int(SLIDE_H * 0.0025), theme["primary"])

        margin = int(SLIDE_W * 0.055)
        mid_t  = footer_t + int(footer_h * 0.22)

        # Brand name — left
        if theme["brand_name"]:
            self._add_textbox(
                slide, margin, mid_t,
                int(SLIDE_W * 0.38), int(footer_h * 0.60),
                theme["brand_name"].upper(), 9, theme["muted"], bold=True,
            )
        # Website URL — centre
        if theme.get("website_url"):
            self._add_textbox(
                slide, int(SLIDE_W * 0.38), mid_t,
                int(SLIDE_W * 0.30), int(footer_h * 0.60),
                theme["website_url"], 8, theme["muted"],
            )
        # Slide counter — right
        self._add_textbox(
            slide,
            SLIDE_W - margin - int(SLIDE_W * 0.12), mid_t,
            int(SLIDE_W * 0.12), int(footer_h * 0.60),
            f"{num:02d} / {total:02d}", 9, theme["muted"], bold=True,
        )

    def _add_rect(self, slide, left, top, width, height, fill_color):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            RECT,
            Emu(int(left)), Emu(int(top)),
            Emu(int(width)), Emu(int(height)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def _add_rounded_rect(self, slide, left, top, width, height, fill_color, corner=0.08):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            ROUND_RECT,
            Emu(int(left)), Emu(int(top)),
            Emu(int(width)), Emu(int(height)),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
        return shape

    def _add_textbox(
        self, slide, left, top, width, height,
        text: str, font_size: float, font_color,
        bold: bool = False, italic: bool = False,
        word_wrap: bool = True,
        vertical_anchor=None,
    ):
        from pptx.util import Emu, Pt
        txBox = slide.shapes.add_textbox(
            Emu(int(left)), Emu(int(top)),
            Emu(int(width)), Emu(int(height)),
        )
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        if vertical_anchor is not None:
            tf.vertical_anchor = vertical_anchor
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(text) if text else ""
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.italic = italic
        return tf

    def _add_slide_header(self, slide, title: str, theme: dict, section_label: str = "") -> None:
        """
        Modern editorial slide header:
          - Full-height left accent strip (primary color) + thin secondary strip
          - Slide title (left, 28pt bold, text_main)
          - Section label displayed as large muted text on the right of header
          - Brand name (top-right, 11pt muted)
          - Full-width muted rule below header zone
        Design style varies per theme["design_variant"].
        """
        variant = theme.get("design_variant", "diagonal")
        margin  = int(SLIDE_W * 0.055)

        if variant in ("diagonal", "bold_contrast"):
            # Full-height left accent strips
            strip_w = int(SLIDE_W * 0.018)
            self._add_rect(slide, 0, 0, strip_w, SLIDE_H, theme["primary"])
            self._add_rect(slide, strip_w, 0, int(strip_w * 0.40), SLIDE_H, theme["secondary"])
        elif variant == "top_band":
            # Top band instead of left strip
            band_h = int(SLIDE_H * 0.014)
            self._add_rect(slide, 0, 0, SLIDE_W, band_h, theme["primary"])
            self._add_rect(slide, 0, band_h, SLIDE_W, int(band_h * 0.45), theme["accent"])
        elif variant == "minimal":
            # Very thin left strip only
            self._add_rect(slide, 0, 0, int(SLIDE_W * 0.008), SLIDE_H, theme["accent"])

        header_top = int(SLIDE_H * 0.038)

        # Slide title
        self._add_textbox(
            slide,
            left=margin, top=header_top,
            width=int(SLIDE_W * 0.60), height=int(SLIDE_H * 0.14),
            text=title,
            font_size=28, font_color=theme["text_main"], bold=True,
        )

        # Brand name (top-right corner)
        self._add_textbox(
            slide,
            left=int(SLIDE_W * 0.82), top=header_top,
            width=int(SLIDE_W * 0.14), height=int(SLIDE_H * 0.10),
            text=theme["brand_name"].upper(),
            font_size=10, font_color=theme["muted"], bold=True,
        )

        # Full-width rule below header
        self._add_rect(
            slide,
            left=0, top=int(SLIDE_H * 0.185),
            width=SLIDE_W, height=int(SLIDE_H * 0.003),
            fill_color=theme["muted"],
        )

    def _add_card(
        self, slide, left, top, width, height, theme,
        label: str = "", header: str = "", body_lines=None,
        accent_color=None, top_border: bool = False,
    ) -> None:
        """
        Draw a card with:
          - Rounded background
          - Colored left border stripe (default) or top border (thicker than before)
          - LABEL text (small, accent-colored)
          - HEADER text (bold, text_main, larger)
          - BODY lines (regular, muted, bulleted)
        """
        body_lines = body_lines or []
        accent = accent_color or theme["primary"]
        # Thicker borders for a more premium look
        border_thickness = int(SLIDE_W * 0.006) if not top_border else int(SLIDE_H * 0.013)

        # Card background
        self._add_rounded_rect(slide, left, top, width, height, theme["card_bg"])

        # Border
        if top_border:
            self._add_rect(slide, left, top, width, border_thickness, accent)
            pad_left = int(SLIDE_W * 0.015)
            pad_top  = border_thickness + int(SLIDE_H * 0.025)
        else:
            self._add_rect(slide, left, top, border_thickness, height, accent)
            pad_left = border_thickness + int(SLIDE_W * 0.015)
            pad_top  = int(SLIDE_H * 0.025)

        content_left  = left + pad_left
        content_w     = width - pad_left - int(SLIDE_W * 0.015)
        y             = top + pad_top

        # Label
        if label:
            lh = int(SLIDE_H * 0.065)
            self._add_textbox(
                slide, content_left, y, content_w, lh,
                label.upper(), 9, accent, bold=True,
            )
            y += lh

        # Header
        if header:
            hh = int(SLIDE_H * 0.11)
            self._add_textbox(
                slide, content_left, y, content_w, hh,
                header, 16, theme["text_main"], bold=True,
            )
            y += hh

        # Body bullets
        if body_lines:
            remaining = max(int(SLIDE_H * 0.06), (top + height) - y - int(SLIDE_H * 0.02))
            body_text = "\n".join(f"• {ln}" for ln in body_lines) if len(body_lines) > 1 else (body_lines[0] if body_lines else "")
            self._add_textbox(
                slide, content_left, y, content_w, remaining,
                body_text, 12, theme["muted"],
            )

    # =========================================================================
    # SLIDE BUILDERS
    # =========================================================================

    def _build_slide_title(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """
        Title slide with visual style driven by theme["design_variant"]:
          - "diagonal"      : Large diagonal freeform split (left panel / right content)
          - "bold_contrast" : Full-bleed left panel + top-right corner triangle
          - "top_band"      : Wide top band + centered content
          - "minimal"       : Clean split with geometric accent + oversized typography
        """
        from pptx.util import Emu
        from pptx.dml.color import RGBColor

        self._set_background(slide, theme["bg"], theme)
        variant = theme.get("design_variant", "diagonal")


        # ── VARIANT: diagonal ─────────────────────────────────────────────────
        if variant == "diagonal":
            # Large diagonal left panel (primary)
            try:
                b = slide.shapes.build_freeform(0, 0)
                b.add_line_to(int(SLIDE_W * 0.50), 0)
                b.add_line_to(int(SLIDE_W * 0.38), SLIDE_H)
                b.add_line_to(0, SLIDE_H)
                p = b.convert_to_shape(); p.fill.solid(); p.fill.fore_color.rgb = theme["primary"]; p.line.fill.background()
            except Exception:
                self._add_rect(slide, 0, 0, int(SLIDE_W * 0.42), SLIDE_H, theme["primary"])
            # Narrower accent overlay
            try:
                b2 = slide.shapes.build_freeform(0, 0)
                b2.add_line_to(int(SLIDE_W * 0.30), 0)
                b2.add_line_to(int(SLIDE_W * 0.22), SLIDE_H)
                b2.add_line_to(0, SLIDE_H)
                p2 = b2.convert_to_shape(); p2.fill.solid(); p2.fill.fore_color.rgb = theme["accent"]; p2.line.fill.background()
            except Exception:
                pass
            # Small top-right triangle (secondary)
            try:
                b3 = slide.shapes.build_freeform(SLIDE_W, 0)
                b3.add_line_to(SLIDE_W, int(SLIDE_H * 0.45))
                b3.add_line_to(int(SLIDE_W * 0.65), 0)
                p3 = b3.convert_to_shape(); p3.fill.solid(); p3.fill.fore_color.rgb = theme["secondary"]; p3.line.fill.background()
            except Exception:
                pass
            content_left = int(SLIDE_W * 0.43)
            content_w    = int(SLIDE_W * 0.50)

        # ── VARIANT: bold_contrast ─────────────────────────────────────────────
        elif variant == "bold_contrast":
            sidebar_w = int(SLIDE_W * 0.22)
            self._add_rect(slide, 0, 0, sidebar_w, SLIDE_H, theme["primary"])
            self._add_rect(slide, sidebar_w, 0, int(sidebar_w * 0.10), SLIDE_H, theme["accent"])
            # Large top-right triangle
            try:
                b = slide.shapes.build_freeform(SLIDE_W, 0)
                b.add_line_to(SLIDE_W, int(SLIDE_H * 0.65))
                b.add_line_to(int(SLIDE_W * 0.55), 0)
                p = b.convert_to_shape(); p.fill.solid(); p.fill.fore_color.rgb = theme["secondary"]; p.line.fill.background()
            except Exception:
                pass
            try:
                b2 = slide.shapes.build_freeform(SLIDE_W, 0)
                b2.add_line_to(SLIDE_W, int(SLIDE_H * 0.40))
                b2.add_line_to(int(SLIDE_W * 0.72), 0)
                p2 = b2.convert_to_shape(); p2.fill.solid(); p2.fill.fore_color.rgb = theme["accent"]; p2.line.fill.background()
            except Exception:
                pass
            content_left = int(SLIDE_W * 0.26)
            content_w    = int(SLIDE_W * 0.46)

        # ── VARIANT: top_band ─────────────────────────────────────────────────
        elif variant == "top_band":
            # Bold horizontal band across top ~40%
            self._add_rect(slide, 0, 0, SLIDE_W, int(SLIDE_H * 0.42), theme["primary"])
            # Accent strip at band bottom
            self._add_rect(slide, 0, int(SLIDE_H * 0.40), SLIDE_W, int(SLIDE_H * 0.025), theme["accent"])
            # Bottom-right corner freeform
            try:
                b = slide.shapes.build_freeform(SLIDE_W, SLIDE_H)
                b.add_line_to(int(SLIDE_W * 0.30), SLIDE_H)
                b.add_line_to(SLIDE_W, int(SLIDE_H * 0.55))
                p = b.convert_to_shape(); p.fill.solid(); p.fill.fore_color.rgb = theme["secondary"]; p.line.fill.background()
            except Exception:
                pass
            content_left = int(SLIDE_W * 0.06)
            content_w    = int(SLIDE_W * 0.58)

        # ── VARIANT: minimal ──────────────────────────────────────────────────
        else:  # "minimal"
            # Thin left strip
            self._add_rect(slide, 0, 0, int(SLIDE_W * 0.012), SLIDE_H, theme["accent"])
            # Large diagonal background element (muted, large, decorative)
            try:
                bg = theme["bg"]
                if theme["dark_theme"]:
                    dec_c = RGBColor(min(255,bg[0]+18), min(255,bg[1]+18), min(255,bg[2]+25))
                else:
                    dec_c = RGBColor(max(0,bg[0]-8), max(0,bg[1]-8), max(0,bg[2]-8))
                b = slide.shapes.build_freeform(SLIDE_W, 0)
                b.add_line_to(SLIDE_W, SLIDE_H)
                b.add_line_to(int(SLIDE_W * 0.38), SLIDE_H)
                p = b.convert_to_shape(); p.fill.solid(); p.fill.fore_color.rgb = dec_c; p.line.fill.background()
            except Exception:
                pass
            # Accent horizontal line across middle
            self._add_rect(slide, int(SLIDE_W * 0.06), int(SLIDE_H * 0.50),
                           int(SLIDE_W * 0.52), int(SLIDE_H * 0.006), theme["primary"])
            content_left = int(SLIDE_W * 0.06)
            content_w    = int(SLIDE_W * 0.58)

        # ── SHARED CONTENT ────────────────────────────────────────────────────
        # Brand name label
        brand_top = int(SLIDE_H * 0.09) if variant == "top_band" else int(SLIDE_H * 0.10)
        brand_color = theme["text_light"] if variant == "top_band" else theme["accent"]
        self._add_textbox(
            slide, content_left, brand_top, content_w, int(SLIDE_H * 0.10),
            theme["brand_name"].upper(), 14, brand_color, bold=True,
        )
        # Accent divider line
        line_top = int(SLIDE_H * 0.225) if variant != "top_band" else int(SLIDE_H * 0.47)
        self._add_rect(slide, content_left, line_top, int(SLIDE_W * 0.44), int(SLIDE_H * 0.006), theme["accent"])

        # Campaign title
        title_top = int(SLIDE_H * 0.26) if variant != "top_band" else int(SLIDE_H * 0.51)
        title_color = theme["text_main"] if variant != "top_band" else theme["text_dark"]
        self._add_textbox(
            slide, content_left, title_top, content_w, int(SLIDE_H * 0.36),
            slide_info.get("title", ""), 38, title_color, bold=True,
        )

        # Subtitle
        subtitle = slide_info.get("subtitle", "")
        if subtitle:
            sub_top = int(SLIDE_H * 0.64) if variant != "top_band" else int(SLIDE_H * 0.70)
            self._add_textbox(
                slide, content_left, sub_top, content_w, int(SLIDE_H * 0.18),
                subtitle, 17, theme["muted"],
            )

        self._add_slide_counter(slide, slide_num, total_slides, theme)

    def _build_slide_company_intro(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """
        Company intro: left 43% = headline + description with left bar,
                       right 52% = 3 stacked KVP cards.
        """
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", "About The Company"), theme, "OVERVIEW")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        margin      = int(SLIDE_W * 0.05)
        content_top = int(SLIDE_H * 0.22)
        content_h   = int(SLIDE_H * 0.72)
        left_w      = int(SLIDE_W * 0.40)

        # Headline
        headline = slide_info.get("headline", "")
        self._add_textbox(
            slide, margin, content_top, left_w, int(SLIDE_H * 0.26),
            headline, 22, theme["text_main"], bold=True,
        )

        # Description with left accent bar
        desc = slide_info.get("description", "")
        if desc:
            bar_top = content_top + int(SLIDE_H * 0.28)
            bar_w   = int(SLIDE_W * 0.006)
            bar_h   = int(SLIDE_H * 0.30)
            self._add_rect(slide, margin, bar_top, bar_w, bar_h, theme["primary"])
            self._add_textbox(
                slide,
                margin + bar_w + int(SLIDE_W * 0.015), bar_top,
                left_w - bar_w - int(SLIDE_W * 0.015), bar_h,
                desc, 13, theme["muted"],
            )

        # Right column: KVP cards
        right_left = int(SLIDE_W * 0.48)
        right_w    = SLIDE_W - right_left - margin
        kvps       = slide_info.get("kvp", [])
        gap        = int(SLIDE_H * 0.018)
        n          = max(len(kvps[:3]), 1)
        card_h     = (content_h - gap * (n - 1)) // n

        accent_colors = [theme["primary"], theme["secondary"], theme["accent"]]
        for i, kvp in enumerate(kvps[:3]):
            card_top = content_top + i * (card_h + gap)
            label    = kvp.get("label", "")
            desc_kv  = kvp.get("description", "")
            accent   = accent_colors[i % 3]

            self._add_rounded_rect(slide, right_left, card_top, right_w, card_h, theme["card_bg"])
            border_w = int(SLIDE_W * 0.005)
            self._add_rect(slide, right_left, card_top, border_w, card_h, accent)

            pad_l = right_left + border_w + int(SLIDE_W * 0.015)
            cw    = right_w - border_w - int(SLIDE_W * 0.025)
            mid   = card_top + int(card_h * 0.15)

            self._add_textbox(slide, pad_l, mid, cw, int(SLIDE_H * 0.07), label, 9, accent, bold=True)
            self._add_textbox(slide, pad_l, mid + int(SLIDE_H * 0.07), cw, int(SLIDE_H * 0.13), desc_kv, 13, theme["text_main"], bold=True)

    def _build_slide_two_by_two(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """2×2 grid of labeled cards."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", ""), theme, "MARKET")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        margin     = int(SLIDE_W * 0.04)
        c_top      = int(SLIDE_H * 0.215)
        c_h        = SLIDE_H - c_top - int(SLIDE_H * 0.05)
        gap_x      = int(SLIDE_W * 0.022)
        gap_y      = int(SLIDE_H * 0.022)
        total_w    = SLIDE_W - 2 * margin
        card_w     = (total_w - gap_x) // 2
        card_h     = (c_h - gap_y) // 2

        accent_colors = [theme["primary"], theme["secondary"], theme["accent"], theme["secondary"]]
        cards = slide_info.get("cards", [])

        for i, card in enumerate(cards[:4]):
            col  = i % 2
            row  = i // 2
            left = margin + col * (card_w + gap_x)
            top  = c_top + row * (card_h + gap_y)

            self._add_card(
                slide, left, top, card_w, card_h, theme,
                label=card.get("label", ""),
                header=card.get("header", ""),
                body_lines=[card.get("body", "")],
                accent_color=accent_colors[i],
            )

    def _build_slide_single_card(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """Single large centered card for a focused challenge or key message."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", ""), theme, "CHALLENGE")

        margin = int(SLIDE_W * 0.06)
        c_top  = int(SLIDE_H * 0.225)
        c_h    = int(SLIDE_H * 0.68)
        c_w    = SLIDE_W - 2 * margin

        # Card background
        self._add_rounded_rect(slide, margin, c_top, c_w, c_h, theme["card_bg"])

        # Left accent border
        border_w = int(SLIDE_W * 0.007)
        self._add_rect(slide, margin, c_top, border_w, c_h, theme["accent"])

        # Icon placeholder (colored circle)
        icon_sz = int(SLIDE_H * 0.13)
        icon_l  = margin + int(SLIDE_W * 0.045)
        icon_t  = c_top + int((c_h - icon_sz) / 2) - int(SLIDE_H * 0.06)
        self._add_rounded_rect(slide, icon_l, icon_t, icon_sz, icon_sz, theme["secondary"], corner=0.5)

        # Label below icon
        label = slide_info.get("label", "")
        if label:
            self._add_textbox(
                slide, icon_l - int(SLIDE_W * 0.01), icon_t + icon_sz + int(SLIDE_H * 0.02),
                icon_sz + int(SLIDE_W * 0.02), int(SLIDE_H * 0.08),
                label, 10, theme["muted"],
            )

        # Vertical divider
        divider_l = icon_l + icon_sz + int(SLIDE_W * 0.04)
        self._add_rect(
            slide, divider_l, c_top + int(SLIDE_H * 0.08),
            int(SLIDE_W * 0.002), c_h - int(SLIDE_H * 0.16), theme["muted"],
        )

        # Headline + body (right of divider)
        content_l = divider_l + int(SLIDE_W * 0.03)
        content_w = c_w - (content_l - margin) - int(SLIDE_W * 0.03)

        headline = slide_info.get("headline", "")
        self._add_textbox(
            slide, content_l, c_top + int(SLIDE_H * 0.12),
            content_w, int(SLIDE_H * 0.24),
            headline, 24, theme["text_main"], bold=True,
        )

        body = slide_info.get("body", "")
        self._add_textbox(
            slide, content_l, c_top + int(SLIDE_H * 0.38),
            content_w, int(SLIDE_H * 0.26),
            body, 14, theme["muted"],
        )

        # Footer brand watermark
        self._add_textbox(
            slide, margin, int(SLIDE_H * 0.91), SLIDE_W - 2 * margin, int(SLIDE_H * 0.06),
            theme["brand_name"], 9, theme["muted"],
        )
        self._add_slide_counter(slide, slide_num, total_slides, theme)

    def _build_slide_three_col(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """3-column card layout (audience, barriers, etc.)."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", ""), theme, "STRATEGY")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        margin  = int(SLIDE_W * 0.04)
        c_top   = int(SLIDE_H * 0.215)
        c_h     = SLIDE_H - c_top - int(SLIDE_H * 0.05)
        gap     = int(SLIDE_W * 0.018)
        total_w = SLIDE_W - 2 * margin
        card_w  = (total_w - 2 * gap) // 3

        accent_colors = [theme["primary"], theme["secondary"], theme["accent"]]
        columns = slide_info.get("columns", [])

        for i, col in enumerate(columns[:3]):
            card_left = margin + i * (card_w + gap)
            accent    = accent_colors[i % 3]
            label     = col.get("label", "")
            header    = col.get("header", "")
            items     = col.get("items", [])

            # Card background
            self._add_rounded_rect(slide, card_left, c_top, card_w, c_h, theme["card_bg"])

            # Top colored border
            self._add_rect(slide, card_left, c_top, card_w, int(SLIDE_H * 0.008), accent)

            # Icon area
            icon_sz = int(SLIDE_H * 0.09)
            icon_l  = card_left + int(SLIDE_W * 0.015)
            icon_t  = c_top + int(SLIDE_H * 0.035)
            self._add_rounded_rect(slide, icon_l, icon_t, icon_sz, icon_sz, accent, corner=0.3)

            # Header text (right of icon)
            hdr_l   = icon_l + icon_sz + int(SLIDE_W * 0.012)
            hdr_w   = card_w - (hdr_l - card_left) - int(SLIDE_W * 0.01)
            self._add_textbox(
                slide, hdr_l, icon_t, hdr_w, int(SLIDE_H * 0.10),
                header, 14, theme["text_main"], bold=True,
            )

            # Label tag (below header)
            lbl_top = icon_t + int(SLIDE_H * 0.10)
            self._add_textbox(
                slide, hdr_l, lbl_top, hdr_w, int(SLIDE_H * 0.065),
                label, 8, accent, bold=True,
            )

            # Horizontal rule
            rule_top = c_top + int(SLIDE_H * 0.215)
            self._add_rect(
                slide, card_left + int(SLIDE_W * 0.015), rule_top,
                card_w - int(SLIDE_W * 0.03), int(SLIDE_H * 0.003), theme["muted"],
            )

            # Items
            items_top = rule_top + int(SLIDE_H * 0.025)
            avail_h   = (c_top + c_h) - items_top - int(SLIDE_H * 0.02)
            is_buyers = "BUYER" in label.upper()

            if is_buyers:
                item_h = avail_h // max(len(items[:3]), 1)
                for j, item in enumerate(items[:3]):
                    pt = items_top + j * item_h
                    ph = min(item_h - int(SLIDE_H * 0.01), int(SLIDE_H * 0.105))
                    pw = card_w - int(SLIDE_W * 0.03)
                    self._add_rounded_rect(
                        slide, card_left + int(SLIDE_W * 0.015), pt, pw, ph,
                        theme["secondary"], corner=0.3,
                    )
                    self._add_textbox(
                        slide,
                        card_left + int(SLIDE_W * 0.025), pt + int(ph * 0.12),
                        pw - int(SLIDE_W * 0.02), ph,
                        item, 11, theme["text_main"], bold=True,
                    )
            else:
                body_text = "\n".join(f"• {it}" for it in items[:3])
                self._add_textbox(
                    slide, card_left + int(SLIDE_W * 0.015), items_top,
                    card_w - int(SLIDE_W * 0.03), avail_h,
                    body_text, 11, theme["muted"],
                )

    def _apply_circle_crop(self, pic_shape) -> None:
        """Apply ellipse (circular) crop to a Picture shape via XML."""
        try:
            from lxml import etree
            ns_a  = "http://schemas.openxmlformats.org/drawingml/2006/main"
            sp_pr = pic_shape._element.spPr
            for child in list(sp_pr):
                if child.tag == f"{{{ns_a}}}prstGeom":
                    sp_pr.remove(child)
            prst_geom = etree.Element(f"{{{ns_a}}}prstGeom")
            prst_geom.set("prst", "ellipse")
            etree.SubElement(prst_geom, f"{{{ns_a}}}avLst")
            xfrm = sp_pr.find(f"{{{ns_a}}}xfrm")
            if xfrm is not None:
                sp_pr.insert(list(sp_pr).index(xfrm) + 1, prst_geom)
            else:
                sp_pr.append(prst_geom)
        except Exception as e:
            logger.debug(f"Circle crop failed: {e}")

    def _build_slide_persona_detail(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """
        Detailed persona profile slide:
          Left panel  (~30 % width): circular photo · name · role · company · tags · quote
          Right panel (~65 % width): 2×2 grid of attribute cards with icon + title + body
        """
        from pptx.util import Emu

        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", "Target Audience"), theme, "AUDIENCE")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        # Subtitle line (persona category) sits just below the header rule
        subtitle = slide_info.get("subtitle", "")
        if subtitle:
            self._add_textbox(
                slide,
                int(SLIDE_W * 0.05), int(SLIDE_H * 0.135),
                int(SLIDE_W * 0.65), int(SLIDE_H * 0.055),
                subtitle, 12, theme["accent"], bold=True,
            )

        margin    = int(SLIDE_W * 0.04)
        c_top     = int(SLIDE_H * 0.215)
        c_h       = SLIDE_H - c_top - int(SLIDE_H * 0.05)
        panel_gap = int(SLIDE_W * 0.025)
        left_w    = int(SLIDE_W * 0.30)
        right_l   = margin + left_w + panel_gap
        right_w   = SLIDE_W - right_l - margin

        # ── Left panel ──────────────────────────────────────────────────────────
        self._add_rounded_rect(slide, margin, c_top, left_w, c_h, theme["card_bg"])

        name = slide_info.get("name", "")
        role = slide_info.get("role", "")

        # Circular portrait photo
        photo_sz = int(left_w * 0.55)
        photo_l  = margin + (left_w - photo_sz) // 2
        photo_t  = c_top + int(SLIDE_H * 0.04)

        img_stream = self._generate_persona_image(name, role)
        if img_stream:
            pic = slide.shapes.add_picture(
                img_stream,
                Emu(photo_l), Emu(photo_t),
                Emu(photo_sz), Emu(photo_sz),
            )
            self._apply_circle_crop(pic)
        else:
            self._add_rounded_rect(slide, photo_l, photo_t, photo_sz, photo_sz, theme["primary"], corner=0.5)

        # Name
        lx = margin + int(SLIDE_W * 0.012)
        lw = left_w - int(SLIDE_W * 0.024)
        y  = photo_t + photo_sz + int(SLIDE_H * 0.022)
        self._add_textbox(slide, lx, y, lw, int(SLIDE_H * 0.11), name, 22, theme["text_main"], bold=True)
        y += int(SLIDE_H * 0.11)

        # Role (accent)
        self._add_textbox(slide, lx, y, lw, int(SLIDE_H * 0.065), role, 13, theme["accent"], bold=True)
        y += int(SLIDE_H * 0.065)

        # Company
        company = slide_info.get("company", "")
        if company:
            self._add_textbox(slide, lx, y, lw, int(SLIDE_H * 0.055), company, 10, theme["muted"])
            y += int(SLIDE_H * 0.058)

        # Tag pills
        tags = slide_info.get("tags", [])[:3]
        if tags:
            pill_h   = int(SLIDE_H * 0.052)
            pill_gap = int(SLIDE_W * 0.008)
            pill_w   = (lw - pill_gap * (len(tags) - 1)) // len(tags)
            for j, tag in enumerate(tags):
                pl = lx + j * (pill_w + pill_gap)
                self._add_rounded_rect(slide, pl, y, pill_w, pill_h, theme["secondary"], corner=0.3)
                self._add_textbox(
                    slide, pl + int(SLIDE_W * 0.005), y + int(pill_h * 0.1),
                    pill_w - int(SLIDE_W * 0.01), pill_h,
                    tag, 9, theme["text_main"], bold=True,
                )
            y += pill_h + int(SLIDE_H * 0.022)

        # Divider + "Quote" label + quote text — pinned to bottom of left panel
        quote = slide_info.get("quote", "")
        if quote:
            quote_h   = int(SLIDE_H * 0.10)
            label_h   = int(SLIDE_H * 0.04)
            rule_h    = int(SLIDE_H * 0.003)
            pad       = int(SLIDE_H * 0.02)
            divider_t = (c_top + c_h) - quote_h - label_h - rule_h - int(SLIDE_H * 0.015) - pad
            if divider_t > y:  # only draw if content above doesn't overlap
                self._add_rect(slide, lx, divider_t, lw, rule_h, theme["muted"])
                label_t = divider_t + rule_h + int(SLIDE_H * 0.012)
                self._add_textbox(slide, lx, label_t, lw, label_h, "Quote", 9, theme["muted"])
                quote_t = label_t + label_h
                self._add_textbox(slide, lx, quote_t, lw, quote_h,
                                  f'"{quote}"', 10, theme["text_main"], italic=True)

        # ── Right panel: 2×2 card grid ──────────────────────────────────────────
        cards   = slide_info.get("cards", [])
        gap_x   = int(SLIDE_W * 0.022)
        gap_y   = int(SLIDE_H * 0.028)
        card_w  = (right_w - gap_x) // 2
        card_h  = (c_h - gap_y) // 2
        accents = [theme["primary"], theme["secondary"], theme["accent"], theme["secondary"]]

        for k, card in enumerate(cards[:4]):
            col = k % 2
            row = k // 2
            cl  = right_l + col * (card_w + gap_x)
            ct  = c_top  + row * (card_h + gap_y)
            ac  = accents[k]

            self._add_rounded_rect(slide, cl, ct, card_w, card_h, theme["card_bg"])
            self._add_rect(slide, cl, ct, card_w, int(SLIDE_H * 0.008), ac)

            # Icon square
            icon_sz = int(SLIDE_H * 0.085)
            icon_l  = cl + int(SLIDE_W * 0.015)
            icon_t  = ct + int(SLIDE_H * 0.03)
            self._add_rounded_rect(slide, icon_l, icon_t, icon_sz, icon_sz, ac, corner=0.2)

            # Card title (right of icon) — vertically centred with the icon square
            from pptx.enum.text import MSO_ANCHOR
            title_l = icon_l + icon_sz + int(SLIDE_W * 0.012)
            title_w = card_w - icon_sz - int(SLIDE_W * 0.035)
            self._add_textbox(slide, title_l, icon_t, title_w, icon_sz,
                               card.get("label", ""), 13, theme["text_main"], bold=True,
                               vertical_anchor=MSO_ANCHOR.MIDDLE)

            # Horizontal rule below icon row
            rule_t = ct + int(SLIDE_H * 0.175)
            self._add_rect(slide, cl + int(SLIDE_W * 0.015), rule_t,
                           card_w - int(SLIDE_W * 0.03), int(SLIDE_H * 0.003), theme["muted"])

            # Body text
            body_t = rule_t + int(SLIDE_H * 0.02)
            body_h = (ct + card_h) - body_t - int(SLIDE_H * 0.02)
            self._add_textbox(slide, cl + int(SLIDE_W * 0.015), body_t,
                               card_w - int(SLIDE_W * 0.03), body_h,
                               card.get("body", ""), 11, theme["muted"])

    def _generate_persona_image(self, name: str, role: str):
        """
        Generate a realistic portrait photo for a persona using DALL-E 3.
        Returns a BytesIO stream ready for add_picture(), or None on failure.
        """
        import io
        import urllib.request

        try:
            prompt = (
                f"Professional portrait headshot photograph of a person named {name}. "
                f"{role}. "
                "Warm, approachable expression, smart business-casual attire, "
                "neutral soft grey background, natural studio lighting, "
                "photorealistic, high quality, authentic, diverse."
            )
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1024x1024",
                quality="standard",
                n=1,
            )
            image_url = response.data[0].url
            with urllib.request.urlopen(image_url) as r:
                return io.BytesIO(r.read())
        except Exception as e:
            logger.warning(f"Persona image generation failed for '{name}': {e}")
            return None

    def _build_slide_persona(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """
        Persona slide: each card has an AI-generated portrait photo at the top,
        name + bullet items below. Both images are generated in parallel.
        """
        import concurrent.futures
        from pptx.util import Emu

        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", "Target Personas"), theme, "AUDIENCE")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        margin = int(SLIDE_W * 0.04)
        c_top  = int(SLIDE_H * 0.215)
        c_h    = SLIDE_H - c_top - int(SLIDE_H * 0.05)
        gap    = int(SLIDE_W * 0.022)
        col_w  = (SLIDE_W - 2 * margin - gap) // 2

        personas = [
            (slide_info.get("left", {}),  theme["secondary"]),
            (slide_info.get("right", {}), theme["accent"]),
        ]

        # Generate both portrait images in parallel
        with concurrent.futures.ThreadPoolExecutor(max_workers=2) as pool:
            futures = {
                i: pool.submit(
                    self._generate_persona_image,
                    col_info.get("header", ""),
                    col_info.get("items", [""])[0] if col_info.get("items") else "",
                )
                for i, (col_info, _) in enumerate(personas)
            }
            images = {i: f.result() for i, f in futures.items()}

        for i, (col_info, accent) in enumerate(personas):
            if not col_info:
                continue

            cl     = margin + i * (col_w + gap)
            label  = col_info.get("label", f"PERSONA {chr(65 + i)}")
            name   = col_info.get("header", "")
            items  = col_info.get("items", [])

            # Card background + left border
            self._add_rounded_rect(slide, cl, c_top, col_w, c_h, theme["card_bg"])
            border_w = int(SLIDE_W * 0.005)
            self._add_rect(slide, cl, c_top, border_w, c_h, accent)

            # Portrait photo (or coloured circle fallback)
            photo_sz = int(c_h * 0.38)
            photo_l  = cl + (col_w - photo_sz) // 2
            photo_t  = c_top + int(SLIDE_H * 0.025)

            img_stream = images.get(i)
            if img_stream:
                slide.shapes.add_picture(
                    img_stream,
                    Emu(photo_l), Emu(photo_t),
                    Emu(photo_sz), Emu(photo_sz),
                )
            else:
                self._add_rounded_rect(slide, photo_l, photo_t, photo_sz, photo_sz, accent, corner=0.5)

            # Text area below photo
            text_top = photo_t + photo_sz + int(SLIDE_H * 0.022)
            pad_l    = cl + border_w + int(SLIDE_W * 0.018)
            text_w   = col_w - border_w - int(SLIDE_W * 0.028)

            # Label chip
            self._add_textbox(slide, pad_l, text_top, text_w, int(SLIDE_H * 0.055),
                               label, 8, accent, bold=True)
            text_top += int(SLIDE_H * 0.058)

            # Name (bold, larger)
            self._add_textbox(slide, pad_l, text_top, text_w, int(SLIDE_H * 0.085),
                               name, 13, theme["text_main"], bold=True)
            text_top += int(SLIDE_H * 0.088)

            # Bullet items
            avail_h = (c_top + c_h) - text_top - int(SLIDE_H * 0.02)
            n_items = len(items[:5])
            if n_items:
                item_h = avail_h // n_items
                for j, item in enumerate(items[:5]):
                    self._add_textbox(
                        slide, pad_l, text_top + j * item_h, text_w, item_h,
                        f"• {item}", 9, theme["muted"],
                    )

    def _build_slide_two_col(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """Two-column card comparison."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", ""), theme, "COMPETITIVE")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        margin  = int(SLIDE_W * 0.04)
        c_top   = int(SLIDE_H * 0.215)
        c_h     = SLIDE_H - c_top - int(SLIDE_H * 0.05)
        gap     = int(SLIDE_W * 0.022)
        col_w   = (SLIDE_W - 2 * margin - gap) // 2

        left_info  = slide_info.get("left", {})
        right_info = slide_info.get("right", {})

        self._add_card(
            slide, margin, c_top, col_w, c_h, theme,
            label=left_info.get("label", ""),
            header=left_info.get("header", ""),
            body_lines=left_info.get("items", []),
            accent_color=theme["secondary"],
        )
        self._add_card(
            slide, margin + col_w + gap, c_top, col_w, c_h, theme,
            label=right_info.get("label", ""),
            header=right_info.get("header", ""),
            body_lines=right_info.get("items", []),
            accent_color=theme["accent"],
        )

    def _build_slide_content(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """Standard content slide with consistent header + bullet cards."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", ""), theme)
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        margin    = int(SLIDE_W * 0.05)
        c_top     = int(SLIDE_H * 0.22)
        c_h       = SLIDE_H - c_top - int(SLIDE_H * 0.06)
        content_w = SLIDE_W - 2 * margin
        bullets   = slide_info.get("content", [])

        if not bullets:
            return

        gap   = int(SLIDE_H * 0.018)
        n     = min(len(bullets), 4)
        row_h = (c_h - gap * (n - 1)) // n

        for i, bullet in enumerate(bullets[:4]):
            bt     = c_top + i * (row_h + gap)
            pip_sz = int(SLIDE_H * 0.022)
            # Align pip with first line of text (text anchors to top of box)
            pip_t  = bt + int(SLIDE_H * 0.01)
            pip_gap = int(SLIDE_W * 0.018)

            self._add_rect(slide, margin, pip_t, pip_sz, pip_sz, theme["accent"])
            self._add_textbox(
                slide,
                left=margin + pip_sz + pip_gap, top=bt,
                width=content_w - pip_sz - pip_gap, height=row_h,
                text=bullet, font_size=15, font_color=theme["text_main"],
            )

        # Footer
        self._add_textbox(
            slide, margin, int(SLIDE_H * 0.92), int(SLIDE_W * 0.40), int(SLIDE_H * 0.06),
            theme["brand_name"], 9, theme["muted"],
        )

    def _build_slide_campaign(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """Campaign concept slide with technique badge + numbered tactic markers + step connector."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", ""), theme, "CAMPAIGN")

        margin    = int(SLIDE_W * 0.05)
        c_top     = int(SLIDE_H * 0.215)
        content_w = SLIDE_W - 2 * margin

        # Technique badge (top-right)
        badge_w = int(SLIDE_W * 0.18)
        badge_h = int(SLIDE_H * 0.07)
        badge_l = SLIDE_W - margin - badge_w
        badge_t = int(SLIDE_H * 0.025)
        self._add_rounded_rect(slide, badge_l, badge_t, badge_w, badge_h, theme["accent"], corner=0.3)
        self._add_textbox(
            slide, badge_l + int(SLIDE_W * 0.005), badge_t + int(SLIDE_H * 0.01),
            badge_w - int(SLIDE_W * 0.01), badge_h,
            "CAMPAIGN TECHNIQUE", 8, theme["text_light"], bold=True,
        )

        # Campaign subtitle / concept name
        subtitle = slide_info.get("subtitle", "")
        if subtitle:
            self._add_textbox(
                slide, margin, c_top, int(SLIDE_W * 0.72), int(SLIDE_H * 0.095),
                f'"{subtitle}"', 15, theme["accent"], italic=True,
            )
            c_top += int(SLIDE_H * 0.10)

        # Rule
        self._add_rect(
            slide, margin, c_top, int(SLIDE_W * 0.55), int(SLIDE_H * 0.004), theme["secondary"],
        )
        c_top += int(SLIDE_H * 0.025)

        # Numbered tactic bullets
        bullets   = slide_info.get("content", [])
        n         = min(len(bullets), 3)
        avail_h   = SLIDE_H - c_top - int(SLIDE_H * 0.06)
        row_h     = avail_h // max(n, 1)
        num_sz    = int(SLIDE_H * 0.065)
        text_left = margin + num_sz + int(SLIDE_W * 0.018)
        text_w    = content_w - num_sz - int(SLIDE_W * 0.018)

        # Draw horizontal step connector line before bullets (connecting number badges)
        if n > 1:
            connector_y = c_top + int(row_h * 0.5) + int(num_sz * 0.5)
            connector_x1 = margin + num_sz
            connector_x2 = margin + num_sz + int(SLIDE_W * 0.0015)
            # Line spanning from first to last badge center
            last_num_rt = c_top + (n - 1) * row_h + int((row_h - num_sz) / 2) + int(num_sz * 0.5)
            line_h = last_num_rt - connector_y
            if line_h > 0:
                self._add_rect(
                    slide,
                    margin + int(num_sz * 0.48), connector_y,
                    int(SLIDE_W * 0.003), line_h, theme["secondary"],
                )

        for i, bullet in enumerate(bullets[:3]):
            rt       = c_top + i * row_h
            num_rt   = rt + int((row_h - num_sz) / 2)

            self._add_rounded_rect(slide, margin, num_rt, num_sz, num_sz, theme["secondary"], corner=0.3)
            self._add_textbox(
                slide, margin, num_rt, num_sz, num_sz,
                str(i + 1), 12, theme["text_light"], bold=True,
            )
            self._add_textbox(
                slide, text_left, rt, text_w, row_h,
                bullet, 14, theme["text_main"],
            )

        self._add_textbox(
            slide, margin, int(SLIDE_H * 0.92), int(SLIDE_W * 0.40), int(SLIDE_H * 0.06),
            theme["brand_name"], 9, theme["muted"],
        )
        self._add_slide_counter(slide, slide_num, total_slides, theme)

    def _build_slide_campaign_examples(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """
        Two real company campaign examples side by side.
        Each card: company + technique badge + strategy name + 3 checkmark items.
        """
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", "Campaign Inspiration"), theme, "CAMPAIGN")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        margin  = int(SLIDE_W * 0.04)
        c_top   = int(SLIDE_H * 0.215)
        c_h     = SLIDE_H - c_top - int(SLIDE_H * 0.05)
        gap     = int(SLIDE_W * 0.025)
        col_w   = (SLIDE_W - 2 * margin - gap) // 2

        examples      = slide_info.get("examples", [])
        accent_colors = [theme["primary"], theme["secondary"]]

        for i, ex in enumerate(examples[:2]):
            cl     = margin + i * (col_w + gap)
            accent = accent_colors[i % 2]

            # Card background + top border
            self._add_rounded_rect(slide, cl, c_top, col_w, c_h, theme["card_bg"])
            self._add_rect(slide, cl, c_top, col_w, int(SLIDE_H * 0.008), accent)

            pad = int(SLIDE_W * 0.015)

            # Company name
            self._add_textbox(
                slide, cl + pad, c_top + int(SLIDE_H * 0.025),
                int(col_w * 0.50), int(SLIDE_H * 0.11),
                ex.get("company", ""), 20, theme["text_main"], bold=True,
            )

            # Technique badge (top-right of card)
            bdg_w = int(col_w * 0.44)
            bdg_h = int(SLIDE_H * 0.07)
            bdg_l = cl + col_w - bdg_w - pad
            bdg_t = c_top + int(SLIDE_H * 0.025)
            self._add_rounded_rect(slide, bdg_l, bdg_t, bdg_w, bdg_h, theme["secondary"], corner=0.3)
            self._add_textbox(
                slide, bdg_l + int(SLIDE_W * 0.006), bdg_t + int(SLIDE_H * 0.01),
                bdg_w - int(SLIDE_W * 0.01), bdg_h,
                ex.get("technique", ""), 9, theme["text_light"], bold=True,
            )

            # Strategy theme (accent color)
            self._add_textbox(
                slide, cl + pad, c_top + int(SLIDE_H * 0.15),
                col_w - 2 * pad, int(SLIDE_H * 0.10),
                ex.get("strategy", ""), 13, accent, bold=True,
            )

            # Rule
            rule_top = c_top + int(SLIDE_H * 0.26)
            self._add_rect(
                slide, cl + pad, rule_top,
                col_w - 2 * pad, int(SLIDE_H * 0.003), theme["muted"],
            )

            # Checkmark items
            items    = ex.get("items", [])
            it_top   = rule_top + int(SLIDE_H * 0.03)
            avail    = (c_top + c_h) - it_top - int(SLIDE_H * 0.02)
            item_h   = avail // max(len(items[:3]), 1)
            ck_sz    = int(SLIDE_H * 0.04)

            for j, item in enumerate(items[:3]):
                it_y = it_top + j * item_h
                ct_y = it_y + int((item_h - ck_sz) / 2)

                # Check box
                self._add_rounded_rect(slide, cl + pad, ct_y, ck_sz, ck_sz, accent, corner=0.2)
                self._add_textbox(
                    slide, cl + pad, ct_y, ck_sz, ck_sz,
                    "✓", 10, theme["text_light"], bold=True,
                )
                # Item text
                tx_l = cl + pad + ck_sz + int(SLIDE_W * 0.012)
                tx_w = col_w - pad - ck_sz - int(SLIDE_W * 0.025)
                self._add_textbox(
                    slide, tx_l, it_y, tx_w, item_h,
                    item, 12, theme["text_main"], bold=True,
                )

    def _build_slide_hooks(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """Marketing Hooks slide: alternating colored callout boxes + oversized quote decoration."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", "Marketing Hooks"), theme, "HOOKS")

        margin    = int(SLIDE_W * 0.05)
        c_top     = int(SLIDE_H * 0.215)
        content_w = SLIDE_W - 2 * margin
        hooks     = slide_info.get("content", [])
        n         = min(len(hooks), 5)

        # Oversized decorative opening quote mark behind hooks list
        self._add_textbox(
            slide,
            left=int(SLIDE_W * 0.72), top=int(SLIDE_H * 0.18),
            width=int(SLIDE_W * 0.24), height=int(SLIDE_H * 0.45),
            text="\u201c", font_size=220, font_color=theme["primary"],
            bold=True,
        )

        if n == 0:
            self._add_slide_counter(slide, slide_num, total_slides, theme)
            return

        gap   = int(SLIDE_H * 0.016)
        avail = SLIDE_H - c_top - int(SLIDE_H * 0.06)
        box_h = (avail - gap * (n - 1)) // n

        for i, hook in enumerate(hooks[:n]):
            bt        = c_top + i * (box_h + gap)
            box_color = theme["primary"] if i % 2 == 0 else theme["secondary"]
            edge_w    = int(SLIDE_W * 0.007)

            # Box background
            self._add_rounded_rect(slide, margin, bt, content_w, box_h, box_color)
            # Accent left edge
            self._add_rect(slide, margin, bt, edge_w, box_h, theme["accent"])
            # Hook text
            pad = int(SLIDE_W * 0.02)
            self._add_textbox(
                slide,
                left=margin + edge_w + pad, top=bt + int(box_h * 0.12),
                width=content_w - edge_w - pad, height=int(box_h * 0.76),
                text=hook, font_size=13, font_color=theme["text_light"], bold=True,
            )

        self._add_textbox(
            slide, margin, int(SLIDE_H * 0.93), int(SLIDE_W * 0.40), int(SLIDE_H * 0.06),
            theme["brand_name"], 9, theme["muted"],
        )
        self._add_slide_counter(slide, slide_num, total_slides, theme)

    def _build_slide_kpis(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """
        KPI / Evaluation slide: 3 columns, each with big numbers + labels + progress bars.
        Layout mirrors reference deck Awareness / Engagement / Conversion.
        """
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", "KPI's"), theme, "KPI'S")

        margin  = int(SLIDE_W * 0.04)
        c_top   = int(SLIDE_H * 0.215)
        c_h     = SLIDE_H - c_top - int(SLIDE_H * 0.05)
        gap     = int(SLIDE_W * 0.018)
        total_w = SLIDE_W - 2 * margin
        card_w  = (total_w - 2 * gap) // 3

        accent_colors = [theme["primary"], theme["secondary"], theme["accent"]]
        columns = slide_info.get("columns", [])

        for i, col in enumerate(columns[:3]):
            cl     = margin + i * (card_w + gap)
            accent = accent_colors[i % 3]

            # Card background + top border
            self._add_rounded_rect(slide, cl, c_top, card_w, c_h, theme["card_bg"])
            self._add_rect(slide, cl, c_top, card_w, int(SLIDE_H * 0.008), accent)

            # Icon + column label side by side
            icon_sz = int(SLIDE_H * 0.09)
            icon_l  = cl + int(SLIDE_W * 0.015)
            icon_t  = c_top + int(SLIDE_H * 0.04)
            self._add_rounded_rect(slide, icon_l, icon_t, icon_sz, icon_sz, accent, corner=0.3)

            lbl_l = icon_l + icon_sz + int(SLIDE_W * 0.012)
            lbl_w = card_w - icon_sz - int(SLIDE_W * 0.035)
            self._add_textbox(
                slide, lbl_l, icon_t, lbl_w, int(SLIDE_H * 0.09),
                col.get("label", ""), 15, theme["text_main"], bold=True,
            )
            self._add_textbox(
                slide, lbl_l, icon_t + int(SLIDE_H * 0.085), lbl_w, int(SLIDE_H * 0.065),
                col.get("subtitle", ""), 9, theme["muted"],
            )

            # Rule
            rule_t = c_top + int(SLIDE_H * 0.19)
            self._add_rect(
                slide, cl + int(SLIDE_W * 0.012), rule_t,
                card_w - int(SLIDE_W * 0.024), int(SLIDE_H * 0.003), theme["muted"],
            )

            # Metrics: big number + description
            metrics = col.get("metrics", [])
            m_top   = rule_t + int(SLIDE_H * 0.025)
            avail   = (c_top + c_h) - m_top - int(SLIDE_H * 0.02)
            m_h     = avail // max(len(metrics[:3]), 1)

            for j, m in enumerate(metrics[:3]):
                mt = m_top + j * m_h

                # Big number
                self._add_textbox(
                    slide, cl + int(SLIDE_W * 0.015), mt,
                    card_w - int(SLIDE_W * 0.03), int(SLIDE_H * 0.13),
                    m.get("number", ""), 28, theme["text_main"], bold=True,
                )
                # Progress bar strip (below the number)
                bar_t   = mt + int(SLIDE_H * 0.128)
                bar_w   = card_w - int(SLIDE_W * 0.03)
                bar_h   = int(SLIDE_H * 0.018)
                bar_l   = cl + int(SLIDE_W * 0.015)
                # Background track
                self._add_rounded_rect(slide, bar_l, bar_t, bar_w, bar_h, theme["card_bg"], corner=0.3)
                # Filled portion (~65% as visual indicator)
                self._add_rounded_rect(slide, bar_l, bar_t, int(bar_w * 0.65), bar_h, accent, corner=0.3)
                # Description
                self._add_textbox(
                    slide, cl + int(SLIDE_W * 0.015), bar_t + bar_h + int(SLIDE_H * 0.008),
                    card_w - int(SLIDE_W * 0.03), int(SLIDE_H * 0.065),
                    m.get("desc", ""), 9, theme["muted"],
                )
                # Divider between metrics
                if j < len(metrics) - 1:
                    div_t = mt + m_h - int(SLIDE_H * 0.004)
                    self._add_rect(
                        slide, cl + int(SLIDE_W * 0.015), div_t,
                        card_w - int(SLIDE_W * 0.03), int(SLIDE_H * 0.003), theme["muted"],
                    )

        self._add_slide_counter(slide, slide_num, total_slides, theme)

    def _build_slide_roadmap(self, slide, slide_info: dict, theme: dict, slide_num: int = 1, total_slides: int = 1) -> None:
        """Execution roadmap — 3-phase grid + channel pills."""
        self._set_background(slide, theme["bg"], theme)
        self._add_bg_accent(slide, theme)
        self._add_slide_header(slide, slide_info.get("title", "Execution Roadmap"), theme, "ROADMAP")
        self._add_slide_counter(slide, slide_num, total_slides, theme)

        phases = slide_info.get("phases", [])
        if phases:
            margin  = int(SLIDE_W * 0.04)
            c_top   = int(SLIDE_H * 0.215)
            c_h     = int(SLIDE_H * 0.60)
            gap     = int(SLIDE_W * 0.018)
            total_w = SLIDE_W - 2 * margin
            card_w  = (total_w - 2 * gap) // 3
            accent_colors = [theme["primary"], theme["secondary"], theme["accent"]]
            for i, phase in enumerate(phases[:3]):
                cl     = margin + i * (card_w + gap)
                accent = accent_colors[i % 3]
                self._add_card(
                    slide, cl, c_top, card_w, c_h, theme,
                    label=phase.get("label", f"PHASE {i+1}"),
                    header=phase.get("header", ""),
                    body_lines=phase.get("items", []),
                    accent_color=accent,
                )
        else:
            # Fallback to content bullets if no phases
            self._build_slide_content(slide, slide_info, theme)
            return

        channels = slide_info.get("channels", [])

        pill_top = int(SLIDE_H * 0.84)
        pill_h   = int(SLIDE_H * 0.07)
        pill_gap = int(SLIDE_W * 0.013)
        pill_w   = int(SLIDE_W * 0.125)
        margin   = int(SLIDE_W * 0.05)

        for j, channel in enumerate(channels[:6]):
            pl    = margin + j * (pill_w + pill_gap)
            color = theme["accent"] if j % 2 == 0 else theme["secondary"]
            self._add_rounded_rect(slide, pl, pill_top, pill_w, pill_h, color, corner=0.3)
            pad = int(pill_w * 0.06)
            self._add_textbox(
                slide, pl + pad, pill_top + int(pill_h * 0.12),
                pill_w - pad * 2, int(pill_h * 0.76),
                channel, 10, theme["text_light"], bold=True,
            )


presentation_service = PresentationService()
