import os
import sys
from pathlib import Path
from pptx import Presentation

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))

from app.services.presentation_service import presentation_service

# Mock Slide Data (Output of structure_content)
mock_slides = {
    "slides": [
        {
            "type": "title",
            "title": "EcoFit Brand Strategy",
            "subtitle": "Sustainable Performance"
        },
        {
            "type": "content",
            "title": "The Problem",
            "content": [
                "Plastic waste is destroying oceans.",
                "Gym wear wears out too fast.",
                "Consumers feel guilty."
            ]
        }
    ]
}

def test_pptx():
    print("--- Testing PPTX Generation ---")
    output_file = "test_presentation.pptx"
    
    # Clean up previous run
    if os.path.exists(output_file):
        os.remove(output_file)
        
    path = presentation_service.generate_pptx(mock_slides, output_file)
    
    if path and os.path.exists(path):
        print(f"✓ SUCCESS: File created at {path}")
        print(f"  File Size: {os.path.getsize(path)} bytes")
        
        # Verify logo is present
        prs = Presentation(path)
        logo_found = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'image') and shape.image:
                    logo_found = True
                    break
                elif hasattr(shape, 'text') and 'yellowHEAD' in shape.text:
                    logo_found = True
                    break
            if logo_found:
                break
        
        if logo_found:
            print(f"  ✓ Logo found in presentation")
        else:
            print(f"  ⚠ Warning: Logo not found in presentation")
    else:
        print("✗ FAILURE: File not created.")

if __name__ == "__main__":
    test_pptx()
