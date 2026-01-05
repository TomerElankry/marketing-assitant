"""
Test suite for yellowHEAD logo integration in presentations.
Tests logo detection, addition to slides, and fallback behavior.
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))

from app.services.presentation_service import presentation_service


def test_logo_detection():
    """Test that logo file is detected correctly"""
    print("--- Testing Logo Detection ---")
    
    if presentation_service.logo_path:
        print(f"✓ SUCCESS: Logo detected at {presentation_service.logo_path}")
        assert os.path.exists(presentation_service.logo_path), "Logo file should exist"
        print(f"  File exists: {os.path.exists(presentation_service.logo_path)}")
        print(f"  File size: {os.path.getsize(presentation_service.logo_path)} bytes")
        return True
    else:
        print("⚠ WARNING: No logo file detected (will use text fallback)")
        return False


def test_logo_on_title_slide():
    """Test that logo appears on title slide"""
    print("\n--- Testing Logo on Title Slide ---")
    
    mock_slides = {
        "slides": [
            {
                "type": "title",
                "title": "Test Brand",
                "subtitle": "Prepared by yellowHEAD AI"
            }
        ]
    }
    
    output_file = "test_logo_title.pptx"
    if os.path.exists(output_file):
        os.remove(output_file)
    
    path = presentation_service.generate_pptx(mock_slides, output_file)
    
    if not path or not os.path.exists(path):
        print("✗ FAILURE: PPTX file not created")
        return False
    
    # Open and check the presentation
    prs = Presentation(path)
    assert len(prs.slides) == 1, "Should have 1 slide"
    
    slide = prs.slides[0]
    shapes = slide.shapes
    
    # Check if logo image or text is present
    logo_found = False
    for shape in shapes:
        if hasattr(shape, 'image') and shape.image:  # Image logo
            logo_found = True
            print(f"✓ SUCCESS: Logo image found on title slide")
            break
        elif hasattr(shape, 'text') and 'yellowHEAD' in shape.text:  # Text logo
            logo_found = True
            print(f"✓ SUCCESS: Logo text found on title slide: '{shape.text[:50]}...'")
            break
    
    if not logo_found:
        print("✗ FAILURE: Logo not found on title slide")
        return False
    
    print(f"  Total shapes on slide: {len(shapes)}")
    os.remove(output_file)
    return True


def test_logo_on_all_slide_types():
    """Test that logo appears on all slide types"""
    print("\n--- Testing Logo on All Slide Types ---")
    
    mock_slides = {
        "slides": [
            {"type": "title", "title": "Test Brand", "subtitle": "Subtitle"},
            {"type": "dashboard", "brand_health_score": 78, "executive_summary": "Test", "summary_tag": "TAG", "key_metrics": []},
            {"type": "strategic_overview", "upcoming_moments": [], "current_wins": [], "critical_risks": [], "growth_unlocks": []},
            {"type": "gaps", "title": "Gaps", "gaps": [{"description": "Gap 1"}]},
            {"type": "platform_analysis", "platforms": [{"name": "Discord", "current_usage": "Usage", "feedback": "Feedback"}]},
            {"type": "insights", "title": "Insights", "subtitle": "Sub", "insights": [{"category": "CREATIVE", "impact": "HIGH", "title": "Title", "description": "Desc", "directive": "Directive"}]},
            {"type": "playbook", "title": "Playbook", "subtitle": "Sub", "acquisition_plays": [], "creative_direction": []},
            {"type": "content", "title": "Content", "content": ["Item 1"]}
        ]
    }
    
    output_file = "test_logo_all_slides.pptx"
    if os.path.exists(output_file):
        os.remove(output_file)
    
    path = presentation_service.generate_pptx(mock_slides, output_file)
    
    if not path or not os.path.exists(path):
        print("✗ FAILURE: PPTX file not created")
        return False
    
    prs = Presentation(path)
    assert len(prs.slides) == len(mock_slides["slides"]), f"Should have {len(mock_slides['slides'])} slides"
    
    slide_types = [s.get("type") for s in mock_slides["slides"]]
    all_have_logo = True
    
    for i, slide in enumerate(prs.slides):
        logo_found = False
        for shape in slide.shapes:
            if hasattr(shape, 'image') and shape.image:
                logo_found = True
                break
            elif hasattr(shape, 'text') and 'yellowHEAD' in shape.text:
                logo_found = True
                break
        
        if logo_found:
            print(f"  ✓ Slide {i+1} ({slide_types[i]}): Logo found")
        else:
            print(f"  ✗ Slide {i+1} ({slide_types[i]}): Logo NOT found")
            all_have_logo = False
    
    if all_have_logo:
        print(f"\n✓ SUCCESS: All {len(prs.slides)} slides have logo")
    else:
        print(f"\n✗ FAILURE: Some slides missing logo")
    
    os.remove(output_file)
    return all_have_logo


def test_logo_positioning():
    """Test that logo is positioned correctly (top-left)"""
    print("\n--- Testing Logo Positioning ---")
    
    mock_slides = {
        "slides": [
            {"type": "content", "title": "Test", "content": ["Item"]}
        ]
    }
    
    output_file = "test_logo_position.pptx"
    if os.path.exists(output_file):
        os.remove(output_file)
    
    path = presentation_service.generate_pptx(mock_slides, output_file)
    
    if not path or not os.path.exists(path):
        print("✗ FAILURE: PPTX file not created")
        return False
    
    prs = Presentation(path)
    slide = prs.slides[0]
    
    # Check logo position (should be in top-left: left < 1 inch, top < 1 inch)
    logo_found = False
    for shape in slide.shapes:
        if hasattr(shape, 'image') and shape.image:
            # Image logo
            if shape.left < Inches(1) and shape.top < Inches(1):
                print(f"✓ SUCCESS: Logo image positioned correctly (left: {shape.left}, top: {shape.top})")
                logo_found = True
                break
        elif hasattr(shape, 'text') and 'yellowHEAD' in shape.text:
            # Text logo
            if shape.left < Inches(1) and shape.top < Inches(1):
                print(f"✓ SUCCESS: Logo text positioned correctly (left: {shape.left}, top: {shape.top})")
                logo_found = True
                break
    
    if not logo_found:
        print("✗ FAILURE: Logo not found or incorrectly positioned")
        return False
    
    os.remove(output_file)
    return True


def test_logo_fallback():
    """Test that text fallback works when logo image is not available"""
    print("\n--- Testing Logo Fallback (Text) ---")
    
    # Temporarily remove logo path to test fallback
    original_logo_path = presentation_service.logo_path
    presentation_service.logo_path = None
    
    try:
        mock_slides = {
            "slides": [
                {"type": "content", "title": "Test", "content": ["Item"]}
            ]
        }
        
        output_file = "test_logo_fallback.pptx"
        if os.path.exists(output_file):
            os.remove(output_file)
        
        path = presentation_service.generate_pptx(mock_slides, output_file)
        
        if not path or not os.path.exists(path):
            print("✗ FAILURE: PPTX file not created")
            return False
        
        prs = Presentation(path)
        slide = prs.slides[0]
        
        # Should have text logo
        text_logo_found = False
        for shape in slide.shapes:
            if hasattr(shape, 'text') and 'yellowHEAD' in shape.text:
                text_logo_found = True
                print(f"✓ SUCCESS: Text fallback logo found: '{shape.text[:50]}...'")
                break
        
        if not text_logo_found:
            print("✗ FAILURE: Text fallback logo not found")
            return False
        
        os.remove(output_file)
        return True
        
    finally:
        # Restore original logo path
        presentation_service.logo_path = original_logo_path


def run_all_tests():
    """Run all logo tests"""
    print("=" * 60)
    print("yellowHEAD Logo Integration Tests")
    print("=" * 60)
    
    results = []
    
    results.append(("Logo Detection", test_logo_detection()))
    results.append(("Logo on Title Slide", test_logo_on_title_slide()))
    results.append(("Logo on All Slide Types", test_logo_on_all_slide_types()))
    results.append(("Logo Positioning", test_logo_positioning()))
    results.append(("Logo Fallback", test_logo_fallback()))
    
    print("\n" + "=" * 60)
    print("Test Results Summary")
    print("=" * 60)
    
    passed = sum(1 for _, result in results if result)
    total = len(results)
    
    for test_name, result in results:
        status = "✓ PASS" if result else "✗ FAIL"
        print(f"{status}: {test_name}")
    
    print(f"\nTotal: {passed}/{total} tests passed")
    print("=" * 60)
    
    return passed == total


if __name__ == "__main__":
    success = run_all_tests()
    sys.exit(0 if success else 1)


